"""Private chat-file storage, document extraction, and optional vision analysis."""

import base64
import binascii
import csv
import io
import os
import re
import tempfile
import zipfile
from collections.abc import Sequence
from functools import lru_cache
from pathlib import Path
from xml.etree import ElementTree

import httpx
from docx import Document
from fastapi import HTTPException
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pillow_heif import register_heif_opener
from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import PdfReadError, PdfStreamError
from sqlalchemy.orm import Session

from utto_server.models import Attachment

MAX_ATTACHMENT_BYTES = 50 * 1024 * 1024
MAX_EXTRACTED_CHARACTERS = 24_000
MAX_ATTACHMENT_CONTEXT_CHARACTERS = 36_000
MAX_PDF_PAGES = 80
MAX_SHEET_ROWS = 300
MAX_SLIDES = 80
# Tuned for the user's RTX 4060 Laptop (8 GB): enough detail for photos and document
# screenshots while leaving headroom for Qwen2.5-VL's visual context and video samples.
MAX_VISION_IMAGE_EDGE = 1_024
# Qwen2.5-VL 3B runs locally with a 4K context window. Limit every request to
# three visual inputs so image/GIF/video analysis stays inside that window on
# the RTX 4060 laptop instead of failing after an upload succeeds.
MAX_VISION_PARTS = 3
MAX_VIDEO_FRAMES = 3
MAX_VIDEO_DECODE_FRAMES_PER_SAMPLE = 180
MAX_VISION_RESPONSE_TOKENS = 320
MAX_AUDIO_TRANSCRIPT_CHARACTERS = 24_000
TEXT_MIME_PREFIXES = ("text/",)
TEXT_MIME_TYPES = {
    "application/json",
    "application/javascript",
    "application/xml",
    "application/x-yaml",
    "application/yaml",
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".csv", ".md"}
IMAGE_EXTENSIONS = {".avif", ".bmp", ".gif", ".heic", ".jpeg", ".jpg", ".png", ".webp"}
IMAGE_MIME_PREFIX = "image/"
VIDEO_EXTENSIONS = {".3gp", ".avi", ".m4v", ".mkv", ".mov", ".mp4", ".mpeg", ".mpg", ".webm"}
VIDEO_MIME_PREFIX = "video/"
AUDIO_EXTENSIONS = {
    ".aac",
    ".aif",
    ".aiff",
    ".amr",
    ".flac",
    ".m4a",
    ".m4b",
    ".mp3",
    ".oga",
    ".ogg",
    ".opus",
    ".wav",
    ".weba",
    ".wma",
}
AUDIO_MIME_PREFIX = "audio/"

register_heif_opener()


def _extension(attachment: Attachment) -> str:
    return Path(attachment.filename).suffix.lower()


def is_image_attachment(attachment: Attachment) -> bool:
    return (
        attachment.mime_type.startswith(IMAGE_MIME_PREFIX)
        or _extension(attachment) in IMAGE_EXTENSIONS
    )


def is_video_attachment(attachment: Attachment) -> bool:
    return (
        attachment.mime_type.startswith(VIDEO_MIME_PREFIX)
        or _extension(attachment) in VIDEO_EXTENSIONS
    )


def is_audio_attachment(attachment: Attachment) -> bool:
    return (
        attachment.mime_type.startswith(AUDIO_MIME_PREFIX)
        or _extension(attachment) in AUDIO_EXTENSIONS
    )


def is_audio_bearing_attachment(attachment: Attachment) -> bool:
    """Whether an attachment may contain speech that can be transcribed locally."""
    return is_audio_attachment(attachment) or is_video_attachment(attachment)


def is_text_attachment(attachment: Attachment) -> bool:
    return (
        attachment.mime_type.startswith(TEXT_MIME_PREFIXES)
        or attachment.mime_type in TEXT_MIME_TYPES
        or _extension(attachment) in DOCUMENT_EXTENSIONS
    )


def decode_upload(content_base64: str) -> bytes:
    try:
        content = base64.b64decode(content_base64, validate=True)
    except (binascii.Error, ValueError) as exc:
        raise HTTPException(status_code=422, detail="Invalid file data") from exc
    return validate_upload(content)


def validate_upload(content: bytes) -> bytes:
    if not content:
        raise HTTPException(status_code=422, detail="The file is empty")
    if len(content) > MAX_ATTACHMENT_BYTES:
        raise HTTPException(status_code=413, detail="Files must be 50 MB or smaller")
    return content


def _plain_text(content: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "gb18030", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _pdf_text(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    return "\n".join((page.extract_text() or "") for page in reader.pages[:MAX_PDF_PAGES])


def _docx_text(content: bytes) -> str:
    document = Document(io.BytesIO(content))
    sections = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                sections.append(" | ".join(cells))
    return "\n".join(sections)


def _spreadsheet_text(content: bytes) -> str:
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sections: list[str] = []
    for worksheet in workbook.worksheets[:12]:
        rows: list[str] = []
        for index, row in enumerate(worksheet.iter_rows(values_only=True)):
            if index >= MAX_SHEET_ROWS:
                break
            values = [str(value).strip() if value is not None else "" for value in row]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            sections.append(f"[{worksheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sections)


def _presentation_text(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides[:MAX_SLIDES], start=1):
        text = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if text:
            sections.append(f"[第 {index} 页]\n" + "\n".join(text))
    return "\n\n".join(sections)


def _csv_text(content: bytes) -> str:
    rows = csv.reader(io.StringIO(_plain_text(content)))
    return "\n".join(" | ".join(row) for _, row in zip(range(MAX_SHEET_ROWS), rows))


def _extract_docx_fallback(content: bytes) -> str:
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        document_xml = archive.read("word/document.xml")
    root = ElementTree.fromstring(document_xml)
    return "\n".join(node.text or "" for node in root.iter() if node.tag.endswith("}t"))


def extract_text(attachment: Attachment) -> str:
    """Extract the useful text portion of a document without sending it elsewhere."""
    extension = _extension(attachment)
    try:
        if extension == ".pdf" or attachment.mime_type == "application/pdf":
            text = _pdf_text(attachment.content)
        elif extension == ".docx" or "wordprocessingml" in attachment.mime_type:
            text = _docx_text(attachment.content)
        elif extension in {".xlsx", ".xlsm"} or "spreadsheetml" in attachment.mime_type:
            text = _spreadsheet_text(attachment.content)
        elif extension == ".pptx" or "presentationml" in attachment.mime_type:
            text = _presentation_text(attachment.content)
        elif extension == ".csv" or attachment.mime_type in {"text/csv", "application/csv"}:
            text = _csv_text(attachment.content)
        elif (
            attachment.mime_type.startswith(TEXT_MIME_PREFIXES)
            or attachment.mime_type in TEXT_MIME_TYPES
            or extension == ".md"
        ):
            text = _plain_text(attachment.content)
        else:
            return ""
    except (
        KeyError,
        OSError,
        TypeError,
        ValueError,
        PdfReadError,
        PdfStreamError,
        zipfile.BadZipFile,
        ElementTree.ParseError,
    ):
        if extension == ".docx":
            try:
                text = _extract_docx_fallback(attachment.content)
            except (
                KeyError,
                OSError,
                TypeError,
                ValueError,
                PdfReadError,
                PdfStreamError,
                zipfile.BadZipFile,
                ElementTree.ParseError,
            ):
                return ""
        else:
            return ""

    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text[:MAX_EXTRACTED_CHARACTERS]


def _vision_settings() -> tuple[str, str, str] | None:
    """Read an optional OpenAI-compatible vision endpoint configuration."""
    base_url = os.getenv("UTTO_VISION_BASE_URL", "").strip().rstrip("/")
    model = os.getenv("UTTO_VISION_MODEL", "").strip()
    api_key = os.getenv("UTTO_VISION_API_KEY", "").strip()
    if not base_url or not model:
        return None
    return base_url, model, api_key


def vision_enabled() -> bool:
    """Whether image bytes can be meaningfully interpreted for this deployment."""
    return _vision_settings() is not None


def _audio_settings() -> tuple[str, str, str] | None:
    """Read the optional local speech-to-text model configuration."""
    model_path = os.getenv("UTTO_AUDIO_MODEL", "").strip()
    if not model_path:
        return None
    model = str(Path(model_path).expanduser())
    if not Path(model).is_dir():
        return None
    device = os.getenv("UTTO_AUDIO_DEVICE", "cpu").strip() or "cpu"
    compute_type = os.getenv("UTTO_AUDIO_COMPUTE_TYPE", "int8").strip() or "int8"
    return model, device, compute_type


def audio_enabled() -> bool:
    """Whether audio and video speech can be transcribed for the chat model."""
    return _audio_settings() is not None


@lru_cache(maxsize=2)
def _audio_transcriber(model_name: str, device: str, compute_type: str) -> object:
    """Load one local Whisper model directory and retain it for later attachments."""
    from faster_whisper import WhisperModel

    return WhisperModel(model_name, device=device, compute_type=compute_type)


def _audio_transcription(attachment: Attachment) -> str:
    """Transcribe speech locally and cache the result with the uploaded attachment.

    The original audio is written only to a short-lived server-local temporary file because
    faster-whisper decodes files through PyAV. It never leaves the computer or reaches
    DeepSeek; only the resulting text can be used in the current chat context.
    """
    if not is_audio_bearing_attachment(attachment):
        return ""
    if attachment.audio_transcript_state == "ready":
        return (attachment.audio_transcript or "").strip()
    if attachment.audio_transcript_state == "failed":
        return ""

    settings = _audio_settings()
    if not settings:
        return ""

    suffix = _extension(attachment) or ".audio"
    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            prefix="utto-audio-", suffix=suffix, delete=False
        ) as temporary_file:
            temporary_file.write(attachment.content)
            temporary_path = Path(temporary_file.name)

        transcriber = _audio_transcriber(*settings)
        segments, _ = transcriber.transcribe(
            str(temporary_path),
            beam_size=5,
            vad_filter=True,
            condition_on_previous_text=True,
        )
        transcript = "".join(segment.text for segment in segments).strip()
    except (ImportError, OSError, RuntimeError, TypeError, ValueError):
        attachment.audio_transcript_state = "failed"
        return ""
    finally:
        if temporary_path is not None:
            try:
                temporary_path.unlink(missing_ok=True)
            except OSError:
                pass

    if not transcript:
        attachment.audio_transcript_state = "failed"
        return ""

    attachment.audio_transcript = transcript[:MAX_AUDIO_TRANSCRIPT_CHARACTERS]
    attachment.audio_transcript_state = "ready"
    return attachment.audio_transcript


def _frame_indexes(frame_count: int) -> list[int]:
    if frame_count <= 1:
        return [0]
    if frame_count == 2:
        return [0, 1]
    return [0, frame_count // 2, frame_count - 1]


def _vision_image_parts(image: Image.Image, label: str) -> list[dict[str, object]]:
    frame = image.convert("RGB")
    frame.thumbnail((MAX_VISION_IMAGE_EDGE, MAX_VISION_IMAGE_EDGE))
    encoded = io.BytesIO()
    frame.save(encoded, format="JPEG", quality=82, optimize=True)
    image_base64 = base64.b64encode(encoded.getvalue()).decode("ascii")
    return [
        {"type": "text", "text": f"文件：{label}"},
        {
            "type": "image_url",
            "image_url": {"url": "data:image/jpeg;base64," + image_base64, "detail": "low"},
        },
    ]


def _image_content_parts(attachment: Attachment, remaining: int) -> list[dict[str, object]]:
    """Create compact JPEG image_url parts, including key GIF frames when available."""
    if remaining <= 0:
        return []

    try:
        with Image.open(io.BytesIO(attachment.content)) as image:
            frame_count = int(getattr(image, "n_frames", 1))
            parts: list[dict[str, object]] = []
            for frame_index in _frame_indexes(frame_count)[:remaining]:
                image.seek(frame_index)
                label = attachment.filename
                if frame_count > 1:
                    label += f"（第 {frame_index + 1}/{frame_count} 帧）"
                parts.extend(_vision_image_parts(image, label))
            return parts
    except (OSError, UnidentifiedImageError, ValueError):
        return []


def _video_target_pts(stream: object, limit: int) -> list[int]:
    duration = getattr(stream, "duration", None)
    start_time = getattr(stream, "start_time", None) or 0
    if not isinstance(duration, int) or duration <= 0:
        return [start_time]
    if limit <= 1:
        return [start_time]
    return [start_time + round(duration * index / (limit - 1)) for index in range(limit)]


def _video_content_parts(attachment: Attachment, remaining: int) -> list[dict[str, object]]:
    """Sample a few frames locally so an image VLM can answer short-video questions.

    Qwen2.5-VL through Ollama accepts image content. Sampling preserves visible temporal
    changes without uploading the original video to an external provider.
    """
    if remaining <= 0:
        return []

    try:
        import av
    except ImportError:
        return []

    container = None
    try:
        container = av.open(io.BytesIO(attachment.content), mode="r")
        stream = next((item for item in container.streams if item.type == "video"), None)
        if stream is None:
            return []

        target_pts = _video_target_pts(stream, min(remaining, MAX_VIDEO_FRAMES))
        parts: list[dict[str, object]] = []
        for index, target in enumerate(target_pts, start=1):
            try:
                container.seek(target, stream=stream, any_frame=False, backward=True)
            except (OSError, ValueError):
                if index > 1:
                    break

            chosen_frame = None
            for decoded_count, frame in enumerate(container.decode(stream), start=1):
                chosen_frame = frame
                if (
                    frame.pts is None
                    or frame.pts >= target
                    or decoded_count >= MAX_VIDEO_DECODE_FRAMES_PER_SAMPLE
                ):
                    break
            if chosen_frame is None:
                continue
            parts.extend(
                _vision_image_parts(
                    chosen_frame.to_image(),
                    f"{attachment.filename}（视频采样 {index}/{len(target_pts)}）",
                )
            )
        return parts
    except (OSError, ValueError, EOFError):
        return []
    finally:
        if container is not None:
            container.close()


def _vision_description(attachments: Sequence[Attachment]) -> str:
    """Use a separately configured vision model so the text chat model stays unchanged."""
    settings = _vision_settings()
    visual_attachments = [
        attachment
        for attachment in attachments
        if is_image_attachment(attachment) or is_video_attachment(attachment)
    ]
    if not settings or not visual_attachments:
        return ""

    parts: list[dict[str, object]] = [
        {
            "type": "text",
            "text": (
                "请解析这些用户聊天附件。逐项提取清晰可见的文字，并简洁描述人物、物体、场景；"
                "GIF 和视频请概括可见的关键变化与过程。视频只提供采样画面，不能依据音频推断。"
                "不要猜测身份、关系、意图或情绪；"
                "看不清就明确写无法辨认。用中文，总计不超过 220 个汉字。"
            ),
        }
    ]
    image_count = 0
    for attachment in visual_attachments:
        visual_parts = (
            _image_content_parts(attachment, MAX_VISION_PARTS - image_count)
            if is_image_attachment(attachment)
            else _video_content_parts(attachment, MAX_VISION_PARTS - image_count)
        )
        parts.extend(visual_parts)
        image_count += sum(1 for part in visual_parts if part.get("type") == "image_url")
        if image_count >= MAX_VISION_PARTS:
            break

    if image_count == 0:
        return ""

    base_url, model, api_key = settings
    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    try:
        with httpx.Client(timeout=httpx.Timeout(90.0, connect=10.0)) as client:
            response = client.post(
                f"{base_url}/chat/completions",
                headers=headers,
                json={
                    "model": model,
                    "messages": [
                        {"role": "system", "content": "你是准确、克制的视觉文件解析器。"},
                        {"role": "user", "content": parts},
                    ],
                    "max_tokens": MAX_VISION_RESPONSE_TOKENS,
                    "stream": False,
                },
            )
            response.raise_for_status()
            content = response.json()["choices"][0]["message"]["content"]
    except (httpx.HTTPError, KeyError, IndexError, TypeError, ValueError):
        return ""

    if isinstance(content, str):
        return content.strip()[:MAX_EXTRACTED_CHARACTERS]
    if isinstance(content, list):
        text = "\n".join(
            str(part.get("text", "")).strip()
            for part in content
            if isinstance(part, dict) and isinstance(part.get("text"), str)
        )
        return text.strip()[:MAX_EXTRACTED_CHARACTERS]
    return ""


def attachment_context(
    db: Session,
    relationship_id: str,
    attachment_ids: Sequence[str],
) -> str:
    """Load only the caller's attachments and create the context supplied to the chat model."""
    if not attachment_ids:
        return ""
    attachments = (
        db.query(Attachment)
        .filter(Attachment.relationship_id == relationship_id, Attachment.id.in_(attachment_ids))
        .all()
    )
    if len(attachments) != len(set(attachment_ids)):
        raise HTTPException(status_code=404, detail="An attachment was not found")

    indexed = {attachment.id: attachment for attachment in attachments}
    ordered_attachments = [indexed[attachment_id] for attachment_id in attachment_ids]
    vision_description = _vision_description(ordered_attachments)
    sections: list[str] = []
    total_characters = 0
    for attachment in ordered_attachments:
        text = extract_text(attachment)
        if text:
            remaining = MAX_ATTACHMENT_CONTEXT_CHARACTERS - total_characters
            if remaining <= 0:
                sections.append(f"用户附带文件：{attachment.filename}（正文过长，未继续注入）")
                continue
            text = text[:remaining]
            sections.append(f"用户附带文件：{attachment.filename}\n---\n{text}\n---")
            total_characters += len(text)
            continue
        transcript = _audio_transcription(attachment)
        if transcript:
            remaining = MAX_ATTACHMENT_CONTEXT_CHARACTERS - total_characters
            if remaining <= 0:
                sections.append(f"用户附带音频：{attachment.filename}（转写内容过长，未继续注入）")
                continue
            transcript = transcript[:remaining]
            label = (
                "用户附带视频的音频转写"
                if is_video_attachment(attachment)
                else "用户附带音频转写"
            )
            sections.append(f"{label}：{attachment.filename}\n---\n{transcript}\n---")
            total_characters += len(transcript)
            continue
        if is_image_attachment(attachment):
            sections.append(f"用户附带图片：{attachment.filename}")
        elif is_video_attachment(attachment):
            sections.append(f"用户附带视频：{attachment.filename}")
        elif is_audio_attachment(attachment):
            sections.append(f"用户附带音频：{attachment.filename}")
        else:
            sections.append(f"用户附带文件：{attachment.filename}（无法提取正文）")

    if vision_description:
        sections.append(f"用户附带图片或视频的视觉解析：\n---\n{vision_description}\n---")
    return "\n\n".join(sections)
